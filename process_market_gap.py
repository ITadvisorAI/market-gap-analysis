
import os
import json
import traceback
import requests
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

# === Google Drive Setup via ENV ===
drive_service = None
try:
    SCOPES = ["https://www.googleapis.com/auth/drive"]
    service_account_info = os.getenv("GOOGLE_SERVICE_ACCOUNT_JSON")
    if not service_account_info:
        raise ValueError("Missing GOOGLE_SERVICE_ACCOUNT_JSON environment variable")
    creds = service_account.Credentials.from_service_account_info(
        json.loads(service_account_info), scopes=SCOPES
    )
    drive_service = build("drive", "v3", credentials=creds)
    print("✅ Google Drive client initialized from ENV")
except Exception as e:
    print(f"❌ Google Drive setup failed: {e}")
    traceback.print_exc()

def download_file(url, dest_path):
    try:
        print(f"⬇️ Downloading: {url}")
        response = requests.get(url)
        response.raise_for_status()
        with open(dest_path, "wb") as f:
            f.write(response.content)
        print(f"✅ Downloaded: {dest_path}")
    except Exception as e:
        print(f"❌ Download failed: {e}")
        traceback.print_exc()

def get_or_create_drive_folder(folder_name):
    query = f"name='{folder_name}' and mimeType='application/vnd.google-apps.folder'"
    response = drive_service.files().list(q=query, fields="files(id)").execute()
    folders = response.get("files", [])
    if folders:
        return folders[0]["id"]
    file_metadata = {
        "name": folder_name,
        "mimeType": "application/vnd.google-apps.folder"
    }
    folder = drive_service.files().create(body=file_metadata, fields="id").execute()
    return folder["id"]

def upload_to_drive(file_path, session_id):
    if not drive_service:
        print("⚠️ Google Drive not configured.")
        return None
    if not os.path.exists(file_path):
        print(f"⚠️ File not found: {file_path}")
        return None
    folder_id = get_or_create_drive_folder(session_id)
    file_metadata = {
        "name": os.path.basename(file_path),
        "parents": [folder_id]
    }
    media = MediaFileUpload(file_path, resumable=True)
    uploaded = drive_service.files().create(body=file_metadata, media_body=media, fields="id").execute()
    file_id = uploaded.get("id")
    return f"https://drive.google.com/file/d/{file_id}/view"

def fill_word_template(path, session_id):
    doc = Document()
    doc.add_heading("Market GAP Analysis Report", level=1)
    doc.add_paragraph(f"Session ID: {session_id}")
    doc.add_paragraph("Summary: Modernization opportunities identified for infrastructure.")
    doc.add_paragraph("Recommendations: Replace obsolete hardware with cloud-native alternatives.")
    doc.save(path)
    print(f"📝 Word report saved: {path}")

def fill_ppt_template(path, session_id):
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    title, subtitle = slide.shapes.title, slide.placeholders[1]
    title.text = "Market GAP Summary"
    subtitle.text = f"Session ID: {session_id}"

    slide2 = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide2.shapes.title.text = "Recommendations"
    slide2.placeholders[1].text = "- Migrate Tier 1 workloads to cloud\n- Refresh unsupported hardware\n- Optimize licensing"
    ppt.save(path)
    print(f"📊 PowerPoint report saved: {path}")

def run_market_gap_analysis(session_id, email, files, webhook, session_folder):
    try:
        print(f"🚀 Starting Market GAP analysis for session: {session_id}")
        os.makedirs(session_folder, exist_ok=True)

        for f in files:
            dest_path = os.path.join(session_folder, f["file_name"])
            download_file(f["file_url"], dest_path)

        docx_path = os.path.join(session_folder, "GAP_Market_Report.docx")
        pptx_path = os.path.join(session_folder, "GAP_Market_Summary.pptx")

        fill_word_template(docx_path, session_id)
        fill_ppt_template(pptx_path, session_id)

        docx_url = upload_to_drive(docx_path, session_id)
        pptx_url = upload_to_drive(pptx_path, session_id)

        print(f"📤 Uploaded Word: {docx_url}")
        print(f"📤 Uploaded PPT: {pptx_url}")

    except Exception as e:
        print(f"💥 Market GAP analysis failed: {e}")
        traceback.print_exc()

__all__ = ['run_market_gap_analysis']
